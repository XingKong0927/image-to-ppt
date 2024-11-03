""" 使用训练好的模型，循环处理每一张图片，提取ppt区域并形成pdf文档
python 3.10.0
torch 2.4.0
torchvision 0.19.0
"""
import os
from PIL import Image

from test_maskrcnn import get_trained_mask_rcnn_model, predict_ppt_regions_with_min_bounding_box, save_cropped_ppt_regions


def to_ppt_image(img_dir, cropped_dir):
    """批量处理照片并保存图片到文件夹中"""
    # 加载模型
    model_path = "mask_rcnn_ppt1.pth"
    model = get_trained_mask_rcnn_model(model_path, num_classes=2)

    imgs = sorted(os.listdir(img_dir))  # 遍历图片文件夹中的所有图片文件

    # 处理每张图片
    for count, img_name in enumerate(imgs):
        image_path = os.path.join(img_dir, img_name)
        if not image_path.lower().endswith((".jpg", ".jpeg", ".png")):
            continue  # 跳过非图片文件

        print(f"Processing No.{count+1}: {image_path}")

        # 预测并生成最小外接矩形
        boxes = predict_ppt_regions_with_min_bounding_box(model, image_path, threshold=0.8, epsilon_factor=0.07)

        # 保存裁剪后的PPT区域
        save_cropped_ppt_regions(image_path, boxes, output_dir=cropped_dir)

def image_to_pdf(output_dir, output_pdf_path):
    """保存图片为PDF文档
    
    Arg:
        output_dir: 裁剪后的图片文件夹
        output_pdf_path: 生成的PDF文件路径
    """
    # 检查文件夹是否为空
    cropped_images = list(sorted(os.listdir(output_dir)))
    if cropped_images:
        # 生成完整路径列表
        image_list = [Image.open(os.path.join(output_dir, img)).convert("RGB") for img in cropped_images]
        # 将第一张图片保存为 PDF，其他图片作为追加
        image_list[0].save(output_pdf_path, save_all=True, append_images=image_list[1:])
        print(f"PPT区域PDF已生成: {output_pdf_path}")
    else:
        print(f"未找到任何图片文件在文件夹 {output_dir} 中。")

def images_to_ppt(output_dir, output_ppt_path, slide_size="widescreen"):
    """保存图片为PPT文档，每张图片填充一个页面
    
    Args:
        output_dir: 裁剪后的图片文件夹
        output_ppt_path: 生成的PPT文件路径
        slide_size: 幻灯片比例模式，可选值为 'widescreen'（16:9） 或 'standard'（4:3）
    """
    from pptx import Presentation
    from pptx.util import Inches
    
    # 创建一个新的PPT文档
    prs = Presentation()

    # 设置幻灯片宽高比
    if slide_size == "widescreen":
        prs.slide_width = Inches(13.33)  # 宽屏16:9比例
        prs.slide_height = Inches(7.5)
    elif slide_size == "standard":
        prs.slide_width = Inches(10)     # 标准4:3比例
        prs.slide_height = Inches(7.5)
    else:
        raise ValueError("slide_size 仅支持 'widescreen' 或 'standard'")

    # 遍历图片文件夹中的所有图片
    images = sorted(os.listdir(output_dir))
    for img_name in images:
        img_path = os.path.join(output_dir, img_name)
        if not img_path.lower().endswith((".jpg", ".jpeg", ".png")):
            continue  # 跳过非图片文件

        # 获取图片尺寸
        with Image.open(img_path) as img:
            img_width, img_height = img.size

        # 创建幻灯片并获取幻灯片尺寸
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        for shape in slide.shapes:
            slide.shapes._spTree.remove(shape._element)     # 移除所有形状，确保页面完全空白
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 计算等比例缩放后的图片尺寸
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            # 图片更宽，以宽度为基准等比例缩放
            new_width = slide_width
            new_height = int(slide_width / img_ratio)
        else:
            # 图片更高，以高度为基准等比例缩放
            new_height = slide_height
            new_width = int(slide_height * img_ratio)

        # 将图片居中添加到幻灯片
        left = (slide_width - new_width) // 2
        top = (slide_height - new_height) // 2
        slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

    # 保存PPT文档
    prs.save(output_ppt_path)
    print(f"PPT文档已生成：{output_ppt_path}")


if __name__ == '__main__':
    img_dir = "data\\raw1"    # 待处理照片路径

    # 创建裁剪区域输出文件夹
    cropped_dir = "data\\cropped_ppt1"
    os.makedirs(cropped_dir, exist_ok=True)

    # 应用训练好的maskrcnn模型预测并裁剪图片，并保存到文件夹中
    to_ppt_image(img_dir, cropped_dir)

    # 使用裁剪后的图片生成PDF
    # image_to_pdf(cropped_dir, output_pdf_path="output_ppt.pdf")

    # 使用裁剪后的图片生成PPT
    images_to_ppt(cropped_dir, output_ppt_path="output_ppt1.pptx", slide_size="widescreen")

