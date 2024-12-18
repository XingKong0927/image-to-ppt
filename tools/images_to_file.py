"""将图片保存为文件

"""
import os
from PIL import Image

def images_to_pdf(output_dir, output_pdf_path):
    """保存图片为PDF文档
    
    Arg:
        output_dir: 裁剪后的图片文件夹
        output_pdf_path: 生成的PDF文件路径
    """
    # 检查文件夹是否为空
    cropped_images = list(sorted(os.listdir(output_dir)))
    if cropped_images:
        # 生成完整路径列表
        image_list = [Image.open(os.path.join(output_dir, img)).convert("RGB") for img in cropped_images]
        # 将第一张图片保存为 PDF，其他图片作为追加
        image_list[0].save(output_pdf_path, save_all=True, append_images=image_list[1:])
        print(f"PPT区域PDF已生成: {output_pdf_path}")
    else:
        print(f"未找到任何图片文件在文件夹 {output_dir} 中。")

def images_to_ppt(output_dir, output_ppt_path, slide_size="widescreen"):
    """保存图片为PPT文档，每张图片填充一个页面
    
    Args:
        output_dir: 裁剪后的图片文件夹
        output_ppt_path: 生成的PPT文件路径
        slide_size: 幻灯片比例模式，可选值为 'widescreen'（16:9） 或 'standard'（4:3）
    """
    from pptx import Presentation
    from pptx.util import Inches
    
    # 创建一个新的PPT文档
    prs = Presentation()

    # 设置幻灯片宽高比
    if slide_size == "widescreen":
        prs.slide_width = Inches(13.33)  # 宽屏16:9比例
        prs.slide_height = Inches(7.5)
    elif slide_size == "standard":
        prs.slide_width = Inches(10)     # 标准4:3比例
        prs.slide_height = Inches(7.5)
    else:
        raise ValueError("slide_size 仅支持 'widescreen' 或 'standard'")

    # 遍历图片文件夹中的所有图片
    images = sorted(os.listdir(output_dir))
    for img_name in images:
        img_path = os.path.join(output_dir, img_name)
        if not img_path.lower().endswith((".jpg", ".jpeg", ".png")):
            continue  # 跳过非图片文件

        # 获取图片尺寸
        with Image.open(img_path) as img:
            img_width, img_height = img.size

        # 创建幻灯片并获取幻灯片尺寸
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        for shape in slide.shapes:
            slide.shapes._spTree.remove(shape._element)     # 移除所有形状，确保页面完全空白
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 计算等比例缩放后的图片尺寸
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            # 图片更宽，以宽度为基准等比例缩放
            new_width = slide_width
            new_height = int(slide_width / img_ratio)
        else:
            # 图片更高，以高度为基准等比例缩放
            new_height = slide_height
            new_width = int(slide_height * img_ratio)

        # 将图片居中添加到幻灯片
        left = (slide_width - new_width) // 2
        top = (slide_height - new_height) // 2
        slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)

    # 保存PPT文档
    prs.save(output_ppt_path)
    print(f"PPT文档已生成：{output_ppt_path}")

